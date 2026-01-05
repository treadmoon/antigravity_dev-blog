from pptx import Presentation
from pptx.util import Inches, Pt

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import argparse
    import os


    def add_bullet_slide(prs, title, bullets):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            if i == 0:
                p = body.paragraphs[0]
                p.text = b
            else:
                p = body.add_paragraph()
                p.text = b
            p.level = 0
        return slide


    def add_footer(prs, text):
        if not text:
            return
        slide = prs.slides[0]
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)


    def build_presentation(output_path, author=None, company=None):
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "langchain_demo01 - 前端开发视角"
        slide.placeholders[1].text = "项目简介与开发要点"

        # Overview
        add_bullet_slide(prs, "项目概述", [
            "基于 Next.js (App Router) 的示例项目",
            "演示路由、布局、CSS Modules 与字体配置",
            "目标：前端工程化实践与代码组织示例"
        ])

        # 技术栈
        add_bullet_slide(prs, "技术栈", [
            "Next.js (App Router)",
            "React Server Components + 客户端组件",
            "CSS Modules、next/font、静态资源放 public/",
            "测试：Jest（仓库含示例测试）"
        ])

        # 代码结构
        add_bullet_slide(prs, "代码组织", [
            "src/app: 路由与页面（Server Components 默认为服务端）",
            "@/lib 用于放置共享工具（约定）",
            "components、styles 模块化，遵循项目约定"
        ])

        # 开发流程
        add_bullet_slide(prs, "本地开发与构建", [
            "运行：npm run dev（含 turbopack）",
            "生产构建：npm run build && npm run start",
            "样式与路由遵循 App Router 约定"
        ])

        # 测试与质量
        add_bullet_slide(prs, "测试与代码质量", [
            "仓库包含 Jest 配置与示例测试文件",
            "使用 ESLint（scripts 中含 lint 脚本）",
            "建议：CI 中加入 lint + test 步骤"
        ])

        # 运行与调试
        add_bullet_slide(prs, "运行与调试要点", [
            "优先在服务端获取数据，必要时在客户端组件使用 'use client'",
            "使用 next/link 做内部导航以启用预取",
            "静态资源在 public/ 下，通过绝对路径引用"
        ])

        # 下一步建议
        add_bullet_slide(prs, "下一步建议", [
            "完善文档：贡献指南与本地开发快速开始",
            "为关键页面添加更多测试覆盖",
            "根据需在 CI 中生成并发布演示 PPT 或 HTML"
        ])

        footer_parts = []
        if author:
            footer_parts.append(f"作者: {author}")
        if company:
            footer_parts.append(f"公司: {company}")
        if footer_parts:
            add_footer(prs, ' | '.join(footer_parts))

        out_dir = os.path.dirname(output_path)
        if out_dir and not os.path.exists(out_dir):
            os.makedirs(out_dir, exist_ok=True)

        prs.save(output_path)
        print(f"Saved presentation to {output_path}")


    def parse_args():
        parser = argparse.ArgumentParser(description="Generate project PPTX")
        parser.add_argument("--output", "-o", default="docs/project_presentation.pptx", help="output pptx path")
        parser.add_argument("--author", help="author name to display on slides")
        parser.add_argument("--company", help="company name to display on slides")
        return parser.parse_args()


    if __name__ == "__main__":
        args = parse_args()
        build_presentation(args.output, author=args.author, company=args.company)
